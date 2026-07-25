"""
build_index.py  -  Ingestion & indexing for the Digital Asset Assistant.

This is the "future-state" pipeline in code. For every asset it:
  - reads what is *inside* the file (OCR on images, text on decks/PDFs) so
    filenames no longer have to describe the content,
  - auto-tags asset type + topic,
  - computes a perceptual hash (dHash) to detect near-duplicates,
  - groups versions into an "asset family",
  - makes a thumbnail for images,
and writes everything to index_data/index.json (+ thumbs/).

The shipped repo already contains a pre-built index.json so you can run the
app without any of the libraries below. Re-run this only if you change the data.

Requires (rebuild only):  pillow, pypdf, python-pptx, openpyxl  +  the
`tesseract` OCR binary on your PATH (optional - without it, image text is skipped).
"""

import os
import re
import json
import subprocess
import shutil
import hashlib

import config

# Soft imports so a partial environment still builds *something*.
try:
    from PIL import Image
except Exception:
    Image = None
try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None
try:
    from pptx import Presentation
except Exception:
    Presentation = None
try:
    import openpyxl
except Exception:
    openpyxl = None

IMAGE_EXT = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}


# ---------------------------------------------------------------------------
# Content extraction
# ---------------------------------------------------------------------------
def ocr_image(path):
    """OCR via the tesseract CLI (no python wrapper needed). Returns '' on failure."""
    if not shutil.which("tesseract"):
        return ""
    try:
        out = subprocess.run(
            ["tesseract", path, "-", "--psm", "3"],
            capture_output=True, text=True, timeout=60,
        )
        return re.sub(r"\s+", " ", out.stdout).strip()
    except Exception:
        return ""


def deck_text(path):
    if Presentation is None:
        return "", None, None
    try:
        prs = Presentation(path)
    except Exception:
        return "", None, None
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        chunks.append(t)
    author = created = None
    try:
        cp = prs.core_properties
        author = cp.author or None
        created = cp.created.isoformat() if cp.created else None
    except Exception:
        pass
    return re.sub(r"\s+", " ", " ".join(chunks)).strip(), author, created


def pdf_text(path):
    if PdfReader is None:
        return ""
    try:
        r = PdfReader(path)
        txt = " ".join((p.extract_text() or "") for p in r.pages)
        return re.sub(r"\s+", " ", txt).strip()
    except Exception:
        return ""


# ---------------------------------------------------------------------------
# Perceptual hash (dHash) - pure PIL, no imagehash dependency
# ---------------------------------------------------------------------------
def dhash(path, size=8):
    if Image is None:
        return None
    try:
        img = Image.open(path).convert("L").resize((size + 1, size), Image.LANCZOS)
        px = list(img.getdata())
        w = size + 1
        bits = 0
        i = 0
        for row in range(size):
            for col in range(size):
                left = px[row * w + col]
                right = px[row * w + col + 1]
                bits = (bits << 1) | (1 if left > right else 0)
                i += 1
        return format(bits, "016x")
    except Exception:
        return None


def hamming(a, b):
    if not a or not b:
        return 999
    return bin(int(a, 16) ^ int(b, 16)).count("1")


# ---------------------------------------------------------------------------
# Auto-tagging
# ---------------------------------------------------------------------------
def guess_asset_type(name, folder, ext):
    n = name.lower()
    f = folder.lower()
    if ext == ".pptx":
        return "Slide deck"
    if ext == ".pdf":
        return "PDF document"
    if "logo" in n:
        return "Logo"
    if "banner" in n or "enrollment" in n:
        return "Banner"
    if "infographic" in n or "roadmap" in n:
        return "Infographic"
    if "speaker_card" in n or "speaker" in n:
        return "Promo card"
    if "insta" in n or "post" in n:
        return "Social post"
    if "thumbnail" in n or "thumb" in n:
        return "Thumbnail"
    if "thumbnail" in f:
        return "Thumbnail / photo"
    if ext in IMAGE_EXT:
        return "Image / photo"
    return "File"


def guess_topics(text):
    t = " " + text.lower() + " "
    found = []
    for kw, label in config.TOPIC_KEYWORDS.items():
        if kw in t and label not in found:
            found.append(label)
    return found


def guess_visibility(name, folder):
    """PRIVACY demo: raw personal / unpublished photos are 'restricted'
    (only an admin sees them); published marketing assets are 'team'."""
    n = name.lower()
    personal = ("dhaval.jpg", "daval_at_conference_with_bloomberg_friends.jpg")
    if n in personal:
        return "restricted"
    return "team"


def normalize_family(name):
    """Family key: strip extension, version suffixes, and 'final'/'copy' noise."""
    base = os.path.splitext(name)[0].lower()
    base = re.sub(r"[_\-\s]+(final|final2|finaluse|usethis|copy|v\d+|version\d+)$", "", base)
    base = re.sub(r"[_\-\s]+\d{1,3}$", "", base)          # trailing _07, -12 (chapter/index)
    base = re.sub(r"[^a-z0-9]+", "_", base).strip("_")
    return base


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
def make_thumb(path, asset_id):
    if Image is None:
        return None
    try:
        os.makedirs(config.THUMBS_DIR, exist_ok=True)
        img = Image.open(path).convert("RGB")
        img.thumbnail((360, 360))
        out = os.path.join(config.THUMBS_DIR, asset_id + ".jpg")
        img.save(out, "JPEG", quality=82)
        return asset_id + ".jpg"
    except Exception:
        return None


def index_files():
    assets = []
    if not os.path.isdir(config.ASSETS_DIR):
        print("!! ASSETS_DIR not found:", config.ASSETS_DIR)
        return assets
    for root, _dirs, files in os.walk(config.ASSETS_DIR):
        for fn in sorted(files):
            if fn.startswith("."):
                continue
            path = os.path.join(root, fn)
            ext = os.path.splitext(fn)[1].lower()
            rel = os.path.relpath(path, config.ASSETS_DIR)
            folder = os.path.relpath(root, config.ASSETS_DIR)
            folder = "" if folder == "." else folder
            asset_id = hashlib.md5(rel.encode("utf-8")).hexdigest()[:12]

            ocr = author = created = ""
            body = ""
            if ext in IMAGE_EXT:
                ocr = ocr_image(path)
            elif ext == ".pptx":
                body, author, created = deck_text(path)
            elif ext == ".pdf":
                body = pdf_text(path)

            title = os.path.splitext(fn)[0].replace("_", " ").strip()
            asset_type = guess_asset_type(fn, folder, ext)
            combined = " ".join([title, folder.replace("_", " "), ocr, body])
            topics = guess_topics(combined)
            phash = dhash(path) if ext in IMAGE_EXT else None
            thumb = make_thumb(path, asset_id) if ext in IMAGE_EXT else None

            assets.append({
                "id": asset_id,
                "kind": "file",
                "source": "internal",
                "filename": fn,
                "rel_path": rel.replace("\\", "/"),
                "folder": folder.replace("\\", "/"),
                "title": title,
                "asset_type": asset_type,
                "topics": topics,
                "ocr": ocr,
                "body": body,
                "author": author,
                "created": created,
                "phash": phash,
                "family": normalize_family(fn),
                "visibility": guess_visibility(fn, folder),
                "thumb": thumb,
                "url": None,
            })
            print(f"  indexed {rel:60s} [{asset_type}] ocr={len(ocr)}c body={len(body)}c")
    return assets


def index_public_links():
    out = []
    if openpyxl is None or not os.path.exists(config.PUBLIC_LINKS_XLSX):
        return out
    wb = openpyxl.load_workbook(config.PUBLIC_LINKS_XLSX)
    ws = wb.worksheets[0]
    rows = list(ws.iter_rows(values_only=True))
    for r in rows[1:]:
        if not r or not r[0]:
            continue
        title = str(r[0]).strip()
        url = str(r[1]).strip() if len(r) > 1 and r[1] else ""
        src = str(r[2]).strip() if len(r) > 2 and r[2] else "website"
        aid = hashlib.md5((url or title).encode("utf-8")).hexdigest()[:12]
        topics = guess_topics(title)
        out.append({
            "id": aid, "kind": "link", "source": src,
            "filename": "", "rel_path": "", "folder": "public_links",
            "title": title, "asset_type": f"Public link ({src})",
            "topics": topics, "ocr": "", "body": title,
            "author": None, "created": None, "phash": None,
            "family": normalize_family(title), "visibility": "public",
            "thumb": None, "url": url,
        })
    print(f"  indexed {len(out)} public links")
    return out


def detect_duplicates(assets):
    """Group by family key, then merge families whose images are near-identical."""
    # near-duplicate image detection via perceptual hash
    imgs = [a for a in assets if a.get("phash")]
    for i, a in enumerate(imgs):
        for b in imgs[i + 1:]:
            if a["family"] != b["family"] and hamming(a["phash"], b["phash"]) <= 6:
                b["family"] = a["family"]  # visually identical => same family
    fam = {}
    for a in assets:
        fam.setdefault(a["family"], []).append(a["id"])
    dup_families = {k: v for k, v in fam.items() if len(v) > 1}
    return dup_families


def main():
    os.makedirs(config.INDEX_DIR, exist_ok=True)
    print("Indexing files from:", config.ASSETS_DIR)
    assets = index_files()
    print("Indexing public links from:", config.PUBLIC_LINKS_XLSX)
    assets += index_public_links()
    dup = detect_duplicates(assets)

    payload = {
        "generated_from": config.ASSETS_DIR,
        "count": len(assets),
        "duplicate_families": dup,
        "assets": assets,
    }
    with open(config.INDEX_FILE, "w", encoding="utf-8") as fh:
        json.dump(payload, fh, indent=1, ensure_ascii=False)
    print(f"\nDONE. {len(assets)} assets, {len(dup)} duplicate families.")
    print("Index written to", config.INDEX_FILE)


if __name__ == "__main__":
    main()
